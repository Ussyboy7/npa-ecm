from pathlib import Path
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.section import WD_SECTION
from docx.enum.table import WD_TABLE_ALIGNMENT, WD_CELL_VERTICAL_ALIGNMENT
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt
from pptx.dml.color import RGBColor as PRGB
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path('/Users/macbook')
DOC_OUT = ROOT / 'Documents' / '10_NPA_ECM_Application_Aligned_Revised.docx'
PPT_OUT = ROOT / 'Downloads' / 'NPA ECM_ Institutional Transformation_Application Aligned Revised.pptx'

NAVY = RGBColor(20, 47, 82)
GOLD = RGBColor(183, 137, 44)
DARK = RGBColor(45, 45, 45)
MUTED = RGBColor(90, 90, 90)


def set_cell(cell, text, bold=False, color=DARK, size=9):
    cell.text = ''
    p = cell.paragraphs[0]
    r = p.add_run(text)
    r.bold = bold
    r.font.size = Pt(size)
    r.font.color.rgb = color
    cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER


def add_heading(doc, text, level=1):
    p = doc.add_heading(text, level=level)
    p.style.font.color.rgb = NAVY
    return p


def add_bullets(doc, items, numbered=False):
    for item in items:
        p = doc.add_paragraph(style='List Number' if numbered else 'List Bullet')
        p.add_run(item)


def build_docx():
    doc = Document()
    sec = doc.sections[0]
    sec.top_margin = Inches(0.65)
    sec.bottom_margin = Inches(0.65)
    sec.left_margin = Inches(0.75)
    sec.right_margin = Inches(0.75)
    styles = doc.styles
    styles['Normal'].font.name = 'Aptos'
    styles['Normal'].font.size = Pt(10)
    styles['Normal'].font.color.rgb = DARK
    for name, size in [('Title', 22), ('Heading 1', 15), ('Heading 2', 12)]:
        styles[name].font.name = 'Aptos Display'
        styles[name].font.size = Pt(size)
        styles[name].font.bold = True
        styles[name].font.color.rgb = NAVY

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run('NPA–ECM APPLICATION-ALIGNED DEMO SCRIPT')
    r.bold = True; r.font.size = Pt(20); r.font.color.rgb = NAVY
    p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run('10-Minute Managing Director Demonstration | Revised August 2026')
    r.italic = True; r.font.size = Pt(10); r.font.color.rgb = MUTED
    p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run('Prepared against the implemented NPA-ECM application routes and source modules')
    r.font.size = Pt(9); r.font.color.rgb = GOLD

    add_heading(doc, '1. Demonstration objective', 1)
    doc.add_paragraph('Demonstrate the implemented NPA-ECM correspondence lifecycle across NPA\'s 4 directorates, 28 divisions and 57 departments: register, route through offices, minute or approve, apply and verify an executive seal, search institutional records, monitor performance, and review the audit trail.')
    doc.add_paragraph('The demonstration must show only capabilities currently evidenced in the application. Retention schedules, legal hold, and claims of legal enforceability are not presented as completed capabilities.')

    add_heading(doc, '2. Pre-demo setup', 1)
    add_bullets(doc, [
        'Log in using an approved MD, ED, GM or AGM test account. Do not use production records for the demonstration.',
        'Use a test correspondence with a harmless subject such as “Directive: Port Efficiency KPI Reporting — Q1 2026”.',
        'Confirm the test user has an active office membership and the required role permissions.',
        'Prepare one test attachment and one test approval path. Do not submit an irreversible action without authorization.',
        'Use the application labels: My Inbox, Office Inbox, Register Correspondence, Archives, Search, Verify Seal, Executive Approvals, Executive Dashboard, Performance Analytics and Audit.'
    ])

    add_heading(doc, '3. Ten-minute script', 1)
    table = doc.add_table(rows=1, cols=3)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER; table.style = 'Table Grid'
    for c, t in zip(table.rows[0].cells, ['Time', 'Application action', 'MD-facing narration']): set_cell(c, t, True, RGBColor(255,255,255), 9)
    for c in table.rows[0].cells: c._tc.get_or_add_tcPr().append(__import__('docx').oxml.parse_xml('<w:shd {} w:fill="143052"/>'.format(__import__('docx').oxml.ns.nsdecls('w'))))
    rows = [
        ('0:00–0:45', 'Opening', 'Sir, this demonstration shows the implemented NPA-ECM pipeline for official correspondence, documents, approvals and institutional records. It gives the Authority one controlled place to register work, route it through offices, capture decisions and preserve an audit history.'),
        ('0:45–1:25', 'Orient the sidebar', 'The application is organized around My Workspace, Registry, Cases, Tools, Analytics & Reports and Administration. Today I will use My Inbox, Office Inbox, Search, Verify Seal, Executive Approvals, Analytics and Audit.'),
        ('1:25–2:50', 'Open My Inbox', 'This is the personal work queue: correspondence and documents requiring my action. I can search, filter by status or priority, inspect SLA state, and see pending approvals. The application also supports acting capacity when an authorized officer temporarily occupies an office seat.'),
        ('2:50–4:10', 'Open a correspondence and add a minute', 'The record contains the subject, reference, dates, attachments and workflow context. I can add a minute, choose the action or approval purpose, save a draft, and route the item. The minute and its audit event become part of the official record.'),
        ('4:10–5:35', 'Route to an office', 'The routing control supports office-based routing as well as controlled person-based routing where permitted. I will route this item to the target office in the hierarchy. Office membership, role permissions and acting appointments determine who can process the item.'),
        ('5:35–6:25', 'Approve and apply the executive seal', 'For an authorized approval, the application can apply an executive digital seal using the configured verification controls. The seal carries a serial number, officer and office information, timestamp and document association.'),
        ('6:25–7:10', 'Verify Seal', 'Verify Seal is a public verification route. Entering the serial number returns whether the seal is valid and displays the available identity, office, timestamp and document information. This provides verification evidence; legal enforceability remains subject to NPA policy and legal approval.'),
        ('7:10–8:05', 'Search and Archives', 'Search finds documents, correspondence and cases. The records and archives views support filters by organizational scope, priority and dates, and provide controlled access based on the user’s grade, organization and office memberships.'),
        ('8:05–9:05', 'Executive Dashboard / Performance Analytics', 'The analytics views provide executive and performance visibility, including volumes, turnaround and SLA-related workload. The application also supports overdue indicators, escalation rules and report-oriented views.'),
        ('9:05–9:40', 'Audit', 'The Audit view records actions such as document creation, viewing, routing, approvals, permission changes and related events. Logs can be filtered and exported for compliance review.'),
        ('9:40–10:00', 'Close', 'Sir, NPA-ECM now provides a controlled correspondence lifecycle: register, route, act, approve, seal, verify, search, monitor and audit. The recommended next step is a role-based UAT with Registry and Executive offices, followed by phased rollout after records-retention and governance policies are approved.')
    ]
    for time, action, narration in rows:
        cells = table.add_row().cells
        set_cell(cells[0], time, True, NAVY, 8)
        set_cell(cells[1], action, True, DARK, 8)
        set_cell(cells[2], narration, False, DARK, 8)

    add_heading(doc, '4. Application-verified feature map', 1)
    t = doc.add_table(rows=1, cols=3); t.style = 'Table Grid'; t.alignment = WD_TABLE_ALIGNMENT.CENTER
    for c, txt in zip(t.rows[0].cells, ['Capability', 'Implemented evidence', 'Demo treatment']): set_cell(c, txt, True, RGBColor(255,255,255), 9)
    for c in t.rows[0].cells: c._tc.get_or_add_tcPr().append(__import__('docx').oxml.parse_xml('<w:shd {} w:fill="143052"/>'.format(__import__('docx').oxml.ns.nsdecls('w'))))
    feature_rows = [
        ('Correspondence and minutes', 'Inbox, register, minute modal, minute APIs', 'Show live workflow'),
        ('Office routing and continuity', 'Office IDs, office inboxes, acting appointments', 'Show office routing and explain succession'),
        ('DMS and versioned records', 'DMS routes, upload, document links and versions', 'Show attachment/document link'),
        ('Executive seals', 'Seal generation, serial number, OTP/TOTP controls, verify API', 'Show approval then Verify Seal'),
        ('Analytics and SLA', 'Executive/performance routes, SLA targets and escalation controls', 'Show dashboard and overdue indicators'),
        ('Audit and export', 'Audit page, activity logs, filters and compliance export', 'Show filtered audit event'),
        ('Completion package', 'Backend completion-package service and completion summary document', 'Mention only if the test record visibly generates it'),
        ('Retention and legal hold', 'Documented roadmap, not yet implemented', 'Do not claim as available')
    ]
    for row in feature_rows:
        cells = t.add_row().cells
        for i, val in enumerate(row): set_cell(cells[i], val, i == 0, DARK, 8)

    add_heading(doc, '5. Revised anticipated MD questions', 1)
    qa = [
        ('Why not email and shared drives?', 'They store messages and files, but NPA-ECM adds structured registration, routing, office ownership, approval actions, seals, search and audit events.'),
        ('What happens when an officer changes?', 'Office inboxes and office memberships preserve the work context. Acting Capacity supports controlled temporary succession. The handover process must still follow NPA administrative policy.'),
        ('Are approvals authentic?', 'The application provides executive seal records, serial numbers, timestamps, identity and a public verification route. Legal validity should be confirmed through NPA policy and legal review.'),
        ('Can leadership see bottlenecks?', 'The application provides executive and performance analytics, SLA indicators, overdue workload and escalation-related views.'),
        ('Who can see sensitive documents?', 'Access is scoped through authentication, role permissions, grade level, directorate/division/department and office memberships. This must be validated in UAT with real NPA permission matrices.'),
        ('Is retention compliance complete?', 'Not yet. Archives and soft-delete controls exist, but retention schedules and legal hold require a subsequent implementation and policy approval.'),
        ('What is the next decision?', 'Approve a controlled UAT and phased rollout, beginning with Registry and Executive offices, while formally approving the access, records-retention, legal-hold and seal-governance policies.')
    ]
    for q, a in qa:
        p = doc.add_paragraph(); r = p.add_run(q + ' '); r.bold = True; r.font.color.rgb = NAVY; p.add_run(a)

    add_heading(doc, '6. Demonstration limitations and control notes', 1)
    add_bullets(doc, [
        'Do not describe the audit log as immutable or tamper-proof without an independent security assessment.',
        'Do not describe seals as legally binding without NPA legal and policy confirmation.',
        'Do not claim automated retention schedules, legal hold or disposition workflows as implemented.',
        'Use the application’s actual labels rather than “Documents & Records” or “Administration” as generic replacements.',
        'The source was type-checked and linted successfully; the frontend test suite currently has two date-format expectation failures that should be resolved before production acceptance.'
    ])
    doc.add_paragraph('End of revised application-aligned script.').alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.save(DOC_OUT)


def add_textbox(slide, x, y, w, h, text, size=20, color=RGBColor(255,255,255), bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(PInches(x), PInches(y), PInches(w), PInches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True; tf.margin_left = PInches(0.08); tf.margin_right = PInches(0.08)
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text; r.font.name = 'Aptos'; r.font.size = PPt(size); r.font.bold = bold; r.font.color.rgb = PRGB(color[0], color[1], color[2])
    return box


def add_slide(prs, title, subtitle=None, bullets=None, kicker=None, status=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = PRGB(15, 35, 62)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PInches(0), PInches(0), PInches(13.333), PInches(0.18)); band.fill.solid(); band.fill.fore_color.rgb = PRGB(190, 145, 50); band.line.fill.background()
    if kicker: add_textbox(slide, 0.72, 0.55, 11.8, 0.35, kicker.upper(), 10, (210, 175, 95), True)
    add_textbox(slide, 0.72, 0.95, 11.8, 0.75, title, 27, (255,255,255), True)
    if subtitle: add_textbox(slide, 0.75, 1.8, 11.5, 0.55, subtitle, 13, (205, 215, 225), False)
    if bullets:
        y = 2.55
        for b in bullets:
            add_textbox(slide, 0.95, y, 11.2, 0.5, '• ' + b, 16, (240,245,250), False); y += 0.62
    if status:
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PInches(0.8), PInches(5.9), PInches(11.7), PInches(0.65)); sh.fill.solid(); sh.fill.fore_color.rgb = PRGB(29, 61, 93); sh.line.color.rgb = PRGB(70, 110, 145)
        add_textbox(slide, 1.0, 6.08, 11.3, 0.28, status, 12, (230, 235, 240), False)
    add_textbox(slide, 0.75, 7.05, 8, 0.2, 'Nigerian Ports Authority | NPA–ECM | Application-aligned revision', 8, (160, 180, 195), False)
    return slide


def build_pptx():
    prs = Presentation(); prs.slide_width = PInches(13.333); prs.slide_height = PInches(7.5)
    slide = add_slide(prs, 'NPA–ECM', 'Application-aligned institutional transformation proposal', kicker='Executive proposal', status='This revision distinguishes implemented application capabilities from roadmap and governance-dependent claims.')
    add_textbox(slide, 0.8, 3.05, 11.5, 1.0, 'A controlled pipeline for correspondence, offices, approvals, seals, records, analytics and audit.', 24, (255,255,255), True)
    add_textbox(slide, 0.8, 4.35, 11.2, 0.5, 'Nigerian Ports Authority | August 2026 | Confidential working proposal', 13, (205,215,225), False)

    add_slide(prs, 'Executive overview', 'What the application currently supports', ['Register and manage official correspondence', 'Route work through office inboxes and controlled workflow actions', 'Capture minutes, approvals and attachments', 'Apply and publicly verify executive digital seals', 'Search records, monitor workload and inspect audit events'], kicker='01 | What exists')
    add_slide(prs, 'The institutional risk', 'Scattered files and decisions create operational and accountability risk', ['Paper files, email threads, drafts and personal workspaces fragment context', 'Officer changes can interrupt work when ownership is person-based', 'Unstructured approvals are difficult to verify and retrieve', 'NPA needs a controlled record of the correspondence lifecycle'], kicker='02 | Why')
    add_slide(prs, 'The application response', 'NPA–ECM turns the lifecycle into a governed workflow', ['Register → route → minute → approve → seal → verify → search → monitor → audit', 'Office membership and role permissions determine access and action rights', 'Acting Capacity supports controlled temporary office-seat succession', 'Records remain linked to correspondence, documents and workflow history'], kicker='03 | How')
    add_slide(prs, 'Office-based continuity', 'Implemented office and succession controls', ['Office Inbox separates office work from an individual’s personal queue', 'Correspondence and minutes carry owning/current office context', 'Routing supports target offices and controlled recipients', 'Acting appointments provide time-bound succession when authorized', 'UAT must confirm NPA’s real MD, ED, GM and AGM permission matrix'], kicker='04 | How', status='Accurate claim: continuity is supported by office membership and acting-capacity controls; administrative handover policy still applies.')
    add_slide(prs, 'Correspondence and minutes', 'The working core of NPA–ECM', ['My Inbox shows items requiring the current user’s action', 'Register Correspondence creates controlled records and references', 'Minute actions support action, information, comment and approval purposes', 'Drafts, attachments, forwarding and distribution are part of the workflow', 'Pending approvals and SLA indicators are visible in the work queue'], kicker='05 | What')
    add_slide(prs, 'Documents, archives and search', 'Institutional retrieval with scoped access', ['DMS supports documents, uploads, links, versions and previews', 'Search covers documents, correspondence and cases', 'Archives provide filtered organizational views and controlled retrieval', 'Access is scoped by role, grade, organization and office membership', 'Retention schedules and legal hold are not yet implemented'], kicker='06 | What', status='Do not present full records-retention compliance as complete.')
    add_slide(prs, 'Executive seals and verification', 'Implemented verification evidence', ['Executive approval records can carry a digital seal', 'Seal data includes serial number, officer, office and timestamp', 'OTP/TOTP controls support seal application workflows', 'The Verify Seal route checks a serial number and returns validity information', 'Legal enforceability remains subject to NPA policy and legal review'], kicker='07 | What', status='Use “secure, verifiable executive seal” unless legal validity has been formally approved.')
    add_slide(prs, 'Analytics and audit', 'Operational visibility and accountability', ['Executive and performance analytics routes are implemented', 'Views support workload, turnaround/SLA and overdue indicators', 'Escalation and SLA configuration support operational follow-up', 'Audit records cover document, correspondence, approval and permission events', 'Audit data can be filtered and exported for compliance review'], kicker='08 | What')
    add_slide(prs, 'Completion packages', 'Implemented backend capability; demonstrate only when visible', ['The backend includes a completion-package service and completion-summary documents', 'Correspondence and case records expose completion-package fields', 'The package can preserve the completed record and associated context', 'The demo should show the generated package if the test workflow visibly produces it', 'Do not claim automatic relationship mapping or institutional learning without a visible feature'], kicker='09 | What', status='Application evidence: completion-package service and generated summary document. UI demonstration required for final acceptance.')
    add_slide(prs, 'Governance and implementation boundary', 'NPA scope and what requires policy or further delivery', ['Application scope: 4 directorates, 28 divisions and 57 departments', 'Implemented: workflow, office routing, DMS, seals, verification, analytics and audit', 'Needs UAT: real NPA permissions, hierarchy rules, office assignments and notification behavior', 'Needs governance approval: seal policy, legal status, records classification and access matrix', 'Roadmap: retention schedules, legal hold, disposition and some integration deliverables'], kicker='10 | Governance')
    add_slide(prs, 'Phased path forward', 'A realistic acceptance sequence', ['Phase 1: Registry and Executive offices — correspondence, inboxes, minutes and approvals', 'Phase 2: Directorate and division expansion — office routing, acting capacity and analytics', 'Phase 3: Records governance — retention, legal hold and disposition controls', 'Phase 4: Integrations and production hardening — identity, mail, HRMS/ERP and security assurance', 'Acceptance gates: UAT evidence, security review, policy approval and operational training'], kicker='11 | End')
    add_slide(prs, 'The decision', 'Move from scattered work to a controlled institutional record', ['Approve application-led UAT with representative MD, ED, GM, AGM and Registry roles', 'Approve the office, access, seal and records-governance policies', 'Use the verified demo flow rather than unsupported claims', 'Adopt NPA–ECM as the official channel only after acceptance gates are met'], kicker='12 | Decision', status='NPA–ECM is the memory of the Authority only when the workflow, governance and records policies are all operational.')
    prs.save(PPT_OUT)


if __name__ == '__main__':
    build_docx(); build_pptx(); print(DOC_OUT); print(PPT_OUT)
